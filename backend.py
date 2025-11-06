import os, re, json, zipfile, unicodedata, pandas as pd
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
from typing import TypedDict, Optional
from langchain_core.messages import HumanMessage
from langchain_groq import ChatGroq
from langgraph.graph import StateGraph, END

# =====================================================
# ✅ ENV + LLM SETUP
# =====================================================
load_dotenv()
llm = ChatGroq(
    model_name="openai/gpt-oss-120b",
    api_key=os.getenv("GROQ_API_KEY"),
    temperature=0.3
)

# =====================================================
# ✅ HELPER FUNCTIONS
# =====================================================
def deep_clean_text(text: str) -> str:
    if not text:
        return ""
    text = unicodedata.normalize("NFKC", str(text))
    text = text.replace("\u200b", "").replace("\xa0", " ")
    return re.sub(r"\s+", " ", text).strip()

def extract_from_rels(pptx_path: str):
    mails, phones = set(), set()
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            for name in z.namelist():
                if name.endswith(".rels"):
                    data = z.read(name).decode("utf-8", errors="ignore")
                    mails.update(re.findall(r"mailto:([\w\.-]+@[\w\.-]+)", data))
                    phones.update(re.findall(r"tel:(\+?\d[\d\s-]{7,})", data))
    except Exception:
        pass
    return list(mails), list(phones)

def extract_text_from_pptx(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(deep_clean_text(shape.text))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(deep_clean_text(cell.text))
    mails, phones = extract_from_rels(pptx_path)
    return deep_clean_text(" ".join(texts + mails + phones))

def extract_text_from_docx(docx_path: str) -> str:
    doc = Document(docx_path)
    texts = [deep_clean_text(p.text) for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    texts.append(deep_clean_text(cell.text))
    return " ".join(texts)

def safe_json_extract(response: str) -> dict:
    try:
        start, end = response.find("{"), response.rfind("}")
        return json.loads(response[start:end+1])
    except Exception:
        data = {}
        for line in response.splitlines():
            if ":" in line:
                k, v = line.split(":", 1)
                data[k.strip()] = v.strip()
        return data

def ensure_not_empty(data: dict, required_fields: list) -> dict:
    for key in required_fields:
        if key not in data or not str(data[key]).strip():
            data[key] = "Not specified"
    return data

skills_df = pd.read_csv(r"C:\Users\703437454\Downloads\Master_Skill_Subskill.csv")
internal_skills_dict = {}
for _, row in skills_df.iterrows():
    main_skill = str(row["Main_Skill"]).strip()
    sub_skills = [s.strip() for s in str(row["Sub_Skills"]).split(",") if s.strip()]
    internal_skills_dict[main_skill] = sub_skills

# =====================================================
# ✅ STATE
# =====================================================
class HRState(TypedDict):
    resume: str
    jd: str
    resume_text: Optional[str]
    jd_text: Optional[str]
    name: Optional[str]
    email: Optional[str]
    phone: Optional[str]
    education: Optional[str]
    experience: Optional[int]
    certificates: Optional[str]
    primary_skills: Optional[str]
    secondary_skills: Optional[str]
    additional_skills: Optional[str]
    strengths: Optional[str]
    weaknesses: Optional[str]
    exp_gap: Optional[str]
    match_score: Optional[str]
    matched_points: Optional[str]
    not_matched_points: Optional[str]
    exceedings_points: Optional[str]
    result: Optional[str]
    next_agent: str
    step_completed: list

# =====================================================
# ✅ AGENT 1 — Resume Parser
# =====================================================

def resume_parser(state: HRState):
    resume_text = extract_text_from_pptx(state["resume"])

    # Flatten taxonomy for LLM
    taxonomy_text = "\n".join(
        [f"- {skill}" for main, subs in internal_skills_dict.items() for skill in [main] + subs]
    )

    prompt = f"""
You are an EXPERT resume parser for HR systems.

--- COMPANY SKILL TAXONOMY ---
{taxonomy_text}
--- END TAXONOMY ---

--- RESUME TEXT ---
\"\"\"{resume_text}\"\"\"
--- END RESUME ---

RULES:
1) Primary Skills = skills in the resume that exactly match taxonomy (case-insensitive exact match).
2) Secondary Skills = any skill or technology mentioned in resume that is NOT in taxonomy.
3) Additional Skills = general competencies (leadership, communication, etc.) mentioned in the resume.
4) Extract candidate's name, email, phone, education, experience, and certificates from resume if available.
5) Return **exactly valid JSON ONLY**, no commentary, no guessing beyond what is in the resume.
6) If a category has no entries, return empty list or empty string.

Example:
Resume: 'John Doe, Email: john@example.com, Phone: +91 1234567890, EC2, FastAPI, Leadership'
Taxonomy: 'EC2'
Output:
{{
  "name": "John Doe",
  "email": "john@example.com",
  "phone": "+91 1234567890",
  "education": "",
  "experience": "",
  "certificates": "",
  "primary_skills": ["EC2"],
  "secondary_skills": ["FastAPI"],
  "additional_skills": ["Leadership"]
}}

Produce JSON now.
"""
    try:
        resp = llm.invoke([HumanMessage(content=prompt)]).content
        data = safe_json_extract(resp)
    except Exception as e:
        data = {"error": str(e)}

    # Ensure required fields exist
    required = ["name", "email", "phone", "education", "experience", "certificates",
                "primary_skills", "secondary_skills", "additional_skills"]
    data = ensure_not_empty(data, required)

    return {
        **state,
        **data,
        "next_agent": "JDAnalyzer",
        "step_completed": state["step_completed"] + ["ResumeParser"]
    }

# =====================================================
# ✅ AGENT 2 — JD Analyzer
# =====================================================
def jd_analyzer(state: HRState):
    jd_text = state.get("jd_text") or extract_text_from_docx(state["jd"])
    state["jd_text"] = jd_text
    return {**state, "jd_text": jd_text, "next_agent": "Matcher",
            "step_completed": state["step_completed"] + ["JDAnalyzer"]}

# =====================================================
# ✅ AGENT 3 — Matcher
# =====================================================
def matcher(state: HRState):
    resume_text = state.get("resume_text")
    jd_text = state.get("jd_text")

    prompt = f"""
You are an expert HR evaluator and recruiter.
Compare the candidate’s resume and the job description.
Return only a single match score (0-100).

Resume:
{resume_text}

JD:
{jd_text}

Output JSON:
{{"match_score": <number>}}
"""
    try:
        response = llm.invoke([HumanMessage(content=prompt)]).content
        data = safe_json_extract(response)
    except Exception as e:
        data = {"match_score": "0", "error": str(e)}

    v = str(data.get("match_score", "0")).strip()
    match_score = v if v.isdigit() else "0"

    return {**state, "match_score": match_score, "next_agent": "JDResumeComparator",
            "step_completed": state["step_completed"] + ["Matcher"]}

# =====================================================
# ✅ AGENT 4 — JD-Resume Comparator
# =====================================================
def jd_resume_comparator(state: HRState):
    resume_text = state.get("resume_text")
    jd_text = state.get("jd_text")

    prompt = f"""
Compare the Job Description and Resume below.

Provide JSON with these fields:
{{
  "matched_points": "key strengths or skills that align with JD",
  "not_matched_points": "skills required by JD missing in resume",
  "exceedings_points": "extra or bonus skills in resume not in JD"
}}

Resume:
{resume_text}

JD:
{jd_text}
"""
    try:
        response = llm.invoke([HumanMessage(content=prompt)]).content
        data = safe_json_extract(response)
        data = ensure_not_empty(data, ["matched_points", "not_matched_points", "exceedings_points"])
    except Exception as e:
        data = {"matched_points": "Not specified",
                "not_matched_points": "Not specified",
                "exceedings_points": "Not specified",
                "error": str(e)}

    return {**state, **data, "next_agent": "Insights",
            "step_completed": state["step_completed"] + ["JDResumeComparator"]}

# =====================================================
# ✅ AGENT 5 — Insights (Updated)
# =====================================================
def insights_agent(state: HRState):
    resume_text = state.get("resume_text", "")
    jd_text = state.get("jd_text", "")

    prompt = f"""
Compare the following Resume and Job Description, and provide output in JSON:

{{ 
  "strengths": ["<strength1>", "<strength2>", ...],
  "weaknesses": ["<weakness1>", "<weakness2>", ...] 
}}

Resume:
{resume_text}

JD:
{jd_text}
"""

    try:
        response = llm.invoke([HumanMessage(content=prompt)]).content.strip()
        data = safe_json_extract(response)

        # Ensure fields exist
        strengths = data.get("strengths", [])
        weaknesses = data.get("weaknesses", [])

        return {
            **state,
            "result": response,
            "strengths": " | ".join(strengths) if strengths else "None",
            "weaknesses": " | ".join(weaknesses) if weaknesses else "None",
            "next_agent": "ReportGenerator",
            "step_completed": state["step_completed"] + ["Insights"]
        }

    except Exception as e:
        return {
            **state,
            "result": f"Error: {str(e)}",
            "strengths": "None",
            "weaknesses": "None",
            "next_agent": "ReportGenerator",
            "step_completed": state["step_completed"] + ["Insights (Error)"]
        }

# =====================================================
# ✅ AGENT 6 — Report Generator
# =====================================================
def report_generator(state: HRState):
    return {**state, "next_agent": END,
            "step_completed": state["step_completed"] + ["ReportGenerator"]}

# =====================================================
# ✅ WORKFLOW
# =====================================================
workflow = StateGraph(HRState)
workflow.add_node("ResumeParser", resume_parser)
workflow.add_node("JDAnalyzer", jd_analyzer)
workflow.add_node("Matcher", matcher)
workflow.add_node("JDResumeComparator", jd_resume_comparator)
workflow.add_node("Insights", insights_agent)
workflow.add_node("ReportGenerator", report_generator)
workflow.set_entry_point("ResumeParser")

workflow.add_edge("ResumeParser", "JDAnalyzer")
workflow.add_edge("JDAnalyzer", "Matcher")
workflow.add_edge("Matcher", "JDResumeComparator")
workflow.add_edge("JDResumeComparator", "Insights")
workflow.add_edge("Insights", "ReportGenerator")
workflow.add_edge("ReportGenerator", END)

graph = workflow.compile()

# =====================================================
# ✅ BATCH RUN FUNCTION
# =====================================================
def run_hr_batch(resume_files: list, jd_files: list):
    results = []
    sl_no = 1

    for resume_file in resume_files:
        resume_text = extract_text_from_pptx(resume_file)
        for jd_file in jd_files:
            jd_text = extract_text_from_docx(jd_file)

            state: HRState = {
                "resume": resume_file,
                "jd": jd_file,
                "resume_text": resume_text,
                "jd_text": jd_text,
                "result": "",
                "name": "Not specified",
                "email": "Not specified",
                "phone": "Not specified",
                "education": "Not specified",
                "experience" : "Not specified",
                "certificates": "Not specified",
                "primary_skills": "Not specified",
                "secondary_skills": "Not specified",
                "additional_skills": "Not specified",
                "strengths": "Not specified",
                "weaknesses": "Not specified",
                "matched_points": "Not specified",
                "not_matched_points": "Not specified",
                "exceedings_points": "Not specified",
                "exp_gap": "Not specified",
                "match_score" : "Not specified",
                "next_agent": "ResumeParser",
                "step_completed": []
            }

            final = graph.invoke(state)
            
            def stringify(value):
                if isinstance(value, dict):
                    return ", ".join(f"{kk}: {stringify(vv)}" for kk, vv in value.items())
                elif isinstance(value, (list, tuple)):
                    return ", ".join(stringify(vv) for vv in value)
                elif value is None:
                    return ""
                else:
                    return str(value)

            normalized = {k: stringify(v) for k, v in final.items()}


            results.append({
                "S.No": sl_no,
                "JD_file_name": os.path.basename(jd_file),
                "Resume_file_name": os.path.basename(resume_file),
                "Name": normalized.get("name"),
                "Email": normalized.get("email"),
                "Phone": normalized.get("phone"),
                "Experience": normalized.get("experience"),
                "Primary_skill": normalized.get("primary_skills"),
                "Secondary_skill": normalized.get("secondary_skills"),
                "Additional_skill": normalized.get("additional_skills"),
                "Matched_points": normalized.get("matched_points"),
                "Not_matched_points": normalized.get("not_matched_points"),
                "Exceedings_points": normalized.get("exceedings_points"),
                "Match_Score": int(normalized.get("match_score", "0")) if str(normalized.get("match_score", "0")).isdigit() else 0,
                "Strengths": normalized.get("strengths"),
                "Weaknesses": normalized.get("weaknesses")
            })
            sl_no += 1

    df = pd.DataFrame(results, columns=[
        "JD_file_name","Resume_file_name", "Name", "Email", "Phone", "Experience",
        "Primary_skill", "Secondary_skill", "Matched_points", "Not_matched_points",
        "Exceedings_points", "Match_Score", "Strengths", "Weaknesses"
    ])

    csv_file = "final_match_report.csv"
    df.to_csv(csv_file, index=False)
    return df, csv_file
